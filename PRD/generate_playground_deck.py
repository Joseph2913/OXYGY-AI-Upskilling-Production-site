#!/usr/bin/env python3
"""
Generate OXYGY-branded PowerPoint deck for the Prompt Playground Redesign PRD.
Follows OXYGY skill conventions: Poppins font, brand colors, cover/divider/content layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand constants ──────────────────────────────────────────────────────────
FONT = "Poppins"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

C = {
    "teal":       RGBColor(0x48, 0xBE, 0xAF),
    "dark":       RGBColor(0x2D, 0x37, 0x48),
    "medGray":    RGBColor(0x5C, 0x61, 0x67),
    "panelBg":    RGBColor(0xF2, 0xF4, 0xF6),
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "darkSlate":  RGBColor(0x1A, 0x2B, 0x4A),
    "gold":       RGBColor(0xE8, 0xA8, 0x33),
    "orange":     RGBColor(0xE6, 0x7E, 0x22),
    "redAlert":   RGBColor(0xE5, 0x3E, 0x3E),
    "borderLine": RGBColor(0xD0, 0xD5, 0xDD),
    "lightGray":  RGBColor(0xA0, 0xAE, 0xC0),
    "bodyText":   RGBColor(0x4A, 0x55, 0x68),
}

# Strategy colors from the PRD
STRATEGY_COLORS = {
    "STRUCTURED_BLUEPRINT":     RGBColor(0xC3, 0xD0, 0xF5),  # Soft Lavender
    "CHAIN_OF_THOUGHT":         RGBColor(0xFB, 0xE8, 0xA6),  # Pale Yellow
    "PERSONA_EXPERT_ROLE":      RGBColor(0xA8, 0xF0, 0xE0),  # Mint
    "OUTPUT_FORMAT_SPEC":       RGBColor(0x38, 0xB2, 0xAC),  # Teal
    "CONSTRAINT_FRAMING":       RGBColor(0xFB, 0xCE, 0xB1),  # Soft Peach
    "FEW_SHOT_EXAMPLES":        RGBColor(0xE6, 0xFF, 0xFA),  # Ice Blue
    "ITERATIVE_DECOMPOSITION":  RGBColor(0xC3, 0xD0, 0xF5),  # Soft Lavender
    "TONE_AND_VOICE":           RGBColor(0xFB, 0xE8, 0xA6),  # Pale Yellow
}

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper functions ─────────────────────────────────────────────────────────

def add_shape(slide, shape_type, x, y, w, h, fill=None, border=None, border_width=0.75):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=11, bold=False, color=None,
                 align="left", font_name=FONT, line_spacing=1.15, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    para.alignment = align_map.get(align, PP_ALIGN.LEFT)
    para.space_after = Pt(0)
    para.space_before = Pt(0)
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_rich_text_box(slide, x, y, w, h, segments, font_size=11, align="left",
                      line_spacing=1.3):
    """segments: list of (text, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    para.alignment = align_map.get(align, PP_ALIGN.LEFT)
    for text, bold, color in segments:
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def add_multi_para_box(slide, x, y, w, h, paragraphs, font_size=11, color=None,
                       bold=False, align="left", spacing=6):
    """paragraphs: list of strings, each becomes a paragraph"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for i, text in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        para.space_after = Pt(spacing)
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


# ── COVER SLIDE ──────────────────────────────────────────────────────────────

def make_cover_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark slate background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C["darkSlate"]

    # Teal accent bar at top
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.08, fill=C["teal"])

    # Title
    add_text_box(slide, 0.72, 1.65, 9.0, 1.8,
                 "Prompt Playground\nRedesign",
                 font_size=40, bold=True, color=C["white"], align="left")

    # Subtitle
    add_text_box(slide, 0.72, 3.55, 9.0, 0.45,
                 "Level 1 — Interactive Prompt Engineering Tool",
                 font_size=18, bold=False, color=C["teal"], align="left")

    # Description
    add_text_box(slide, 0.72, 4.2, 8.5, 1.0,
                 "An AI-powered tool that teaches prompt engineering by demonstration — users describe a real task, receive a single optimised prompt, and see the named strategies that built it.",
                 font_size=14, bold=False, color=C["lightGray"], align="left")

    # Date
    add_text_box(slide, 0.72, 6.68, 3.0, 0.26,
                 "March 2026",
                 font_size=11, bold=False, color=C["lightGray"], align="left")

    # Confidentiality
    add_text_box(slide, 9.5, 6.68, 3.0, 0.26,
                 "CONFIDENTIAL",
                 font_size=11, bold=True, color=C["lightGray"], align="right")

    # OXYGY text logo (bottom left area)
    add_text_box(slide, 0.72, 6.2, 3.0, 0.4,
                 "OXYGY",
                 font_size=16, bold=True, color=C["teal"], align="left")


# ── DIVIDER SLIDE ────────────────────────────────────────────────────────────

def make_divider_slide(section_label, section_title, narrative):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Left half: dark slate background panel
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 5.8, 7.5, fill=C["darkSlate"])

    # Teal accent bar on left
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.06, 7.5, fill=C["teal"])

    # Right half: white (default bg)

    # Section label (teal, all-caps)
    add_text_box(slide, 6.2, 1.4, 6.8, 0.28,
                 section_label,
                 font_size=12, bold=True, color=C["teal"], align="left")

    # Section title
    add_text_box(slide, 6.2, 1.75, 6.8, 1.0,
                 section_title,
                 font_size=28, bold=True, color=C["dark"], align="left")

    # Narrative
    add_text_box(slide, 6.2, 2.85, 6.8, 2.8,
                 narrative,
                 font_size=13, bold=False, color=C["medGray"], align="left")

    return slide


# ── CONTENT SLIDE ────────────────────────────────────────────────────────────

def make_content_slide(title, eyebrow, context_banner, footer_takeaway):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Top border line
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.04, fill=C["teal"])

    # Title (ph102 equivalent)
    add_text_box(slide, 0.565, 0.25, 12.202, 0.72,
                 title,
                 font_size=22, bold=True, color=C["dark"], align="left")

    # Eyebrow (ph103 equivalent) - ALL CAPS, teal
    add_text_box(slide, 0.565, 0.97, 12.202, 0.25,
                 eyebrow,
                 font_size=10, bold=True, color=C["teal"], align="left")

    # Context banner (ph104 equivalent) - light bg strip
    banner_shape = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.565, 1.35, 12.202, 0.58,
                             fill=C["panelBg"], border=C["borderLine"], border_width=0.5)
    add_text_box(slide, 0.72, 1.40, 11.9, 0.48,
                 context_banner,
                 font_size=11, bold=False, color=C["dark"], align="left")

    # Footer takeaway (ph105 equivalent)
    footer_shape = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.565, 6.43, 12.202, 0.50,
                             fill=C["panelBg"], border=C["borderLine"], border_width=0.5)
    add_text_box(slide, 0.72, 6.47, 11.9, 0.42,
                 footer_takeaway,
                 font_size=10, bold=False, color=C["medGray"], align="left")

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD THE DECK
# ═══════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Cover ───────────────────────────────────────────────────────────
make_cover_slide()

# ── Slide 2: Section Divider — Overview ──────────────────────────────────────
make_divider_slide(
    "SECTION ONE",
    "Overview & Purpose",
    "The Prompt Playground is the experiential anchor for Level 1. This section establishes what it does, where it sits in the platform, and the key shifts from v1.0 to v2.0 — from a rigid 6-part blueprint to a strategy-aware, open-ended tool."
)

# ── Slide 3: Purpose & Position ──────────────────────────────────────────────
slide = make_content_slide(
    "What the Prompt Playground does and where it lives",
    "OVERVIEW: PURPOSE & POSITION",
    "The Playground teaches prompt engineering by demonstration — users describe a real task, receive a single optimised prompt, and see the prompting strategies combined to build it.",
    "💡 The tool is both a learning instrument and a daily-use utility — internal learners use it for real tasks while absorbing prompting principles unconsciously."
)

# Content zone: two cards side by side
# Card 1: Purpose
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.565, 2.1, 5.95, 4.0,
          fill=C["panelBg"], border=C["borderLine"])
add_text_box(slide, 0.8, 2.25, 5.5, 0.35,
             "PURPOSE", font_size=10, bold=True, color=C["teal"])
add_multi_para_box(slide, 0.8, 2.65, 5.5, 3.2, [
    "• AI-powered tool in the Level 1 section",
    "• Teaches prompt engineering by demonstration",
    "• Single open input → strategy-aware output",
    "• Replaces the two-mode mechanic (Enhance / Build)",
    "• Replaces the mandatory 6-part blueprint output",
    "• Shows that great prompting is a craft with multiple techniques — not a fixed formula"
], font_size=11, color=C["dark"])

# Card 2: Position
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.82, 2.1, 5.95, 4.0,
          fill=C["panelBg"], border=C["borderLine"])
add_text_box(slide, 7.05, 2.25, 5.5, 0.35,
             "POSITION IN SITE", font_size=10, bold=True, color=C["teal"])
add_multi_para_box(slide, 7.05, 2.65, 5.5, 3.2, [
    "• Accessed from Level 1 via 'Try the Prompt Playground' CTA",
    "• Lives at /playground route",
    "• ← Back to Level 1 breadcrumb top-left",
    "• Referenced from Agent Builder (Level 2)",
    "",
    "Two audiences:",
    "  Internal learners → real daily work tool",
    "  External showcase → evaluate Oxygy's depth"
], font_size=11, color=C["dark"])


# ── Slide 4: Key Shifts v1.0 → v2.0 ─────────────────────────────────────────
slide = make_content_slide(
    "Key shifts from v1.0 to v2.0",
    "OVERVIEW: VERSION COMPARISON",
    "The redesign moves from a rigid, single-format output to a flexible, strategy-aware engine that adapts to each task type.",
    "💡 v2.0 positions prompting as a craft with named techniques and judgment calls — not a formula to follow."
)

# Table-like layout: 4 comparison rows
row_data = [
    ("Two mode cards\n(Enhance / Build from Scratch)", "Single open textarea\n— no mode selection"),
    ("6-part Blueprint as\nmandatory output format", "Strategy-aware output:\ncombined prompt + strategy attribution"),
    ("Output always structured\ninto 6 sections", "Output format adapts\nto the task type"),
    ("Educational layer implicit\n(learn by observing structure)", "Educational layer explicit\n(named strategies, practitioner rationale)"),
]

# Column headers
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.565, 2.1, 5.95, 0.45,
          fill=C["dark"])
add_text_box(slide, 0.565, 2.13, 5.95, 0.4,
             "v1.0 (Previous)", font_size=12, bold=True, color=C["white"], align="center")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.82, 2.1, 5.95, 0.45,
          fill=C["teal"])
add_text_box(slide, 6.82, 2.13, 5.95, 0.4,
             "v2.0 (Redesign)", font_size=12, bold=True, color=C["white"], align="center")

y_start = 2.65
row_h = 0.9
for i, (old, new) in enumerate(row_data):
    y = y_start + i * (row_h + 0.1)
    bg = C["white"] if i % 2 == 0 else C["panelBg"]
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.565, y, 5.95, row_h, fill=bg, border=C["borderLine"])
    add_text_box(slide, 0.72, y + 0.1, 5.6, row_h - 0.2, old,
                 font_size=11, color=C["medGray"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.82, y, 5.95, row_h, fill=bg, border=C["borderLine"])
    add_text_box(slide, 7.0, y + 0.1, 5.6, row_h - 0.2, new,
                 font_size=11, color=C["dark"], bold=True)


# ── Slide 5: Section Divider — Strategies ────────────────────────────────────
make_divider_slide(
    "SECTION TWO",
    "The Eight Prompting Strategies",
    "These are the canonical strategies the AI selects from. Each is defined with clear when-to-use and when-not-to-use guidance. Together they form a strategy-aware engine — not a fixed formula."
)


# ── Slide 6: Strategies Overview (1-4) ───────────────────────────────────────
slide = make_content_slide(
    "Eight prompting strategies — the AI's toolkit (1 of 2)",
    "THE STRATEGY ENGINE: STRATEGIES 1–4",
    "Each strategy is a named technique the AI selects from based on the task type. Learners absorb these through repeated exposure — not instruction.",
    "💡 Persona appears in almost every output, reinforcing that 'who is speaking' is the highest-leverage variable in prompting."
)

strategies_1 = [
    ("🏗️", "Structured\nBlueprint", "Breaks prompt into named sections: Role, Context, Task, Format, Steps, Quality Checks.",
     STRATEGY_COLORS["STRUCTURED_BLUEPRINT"]),
    ("🧠", "Chain-of-\nThought", "Instructs AI to reason step-by-step before concluding. Dramatically improves analytical output.",
     STRATEGY_COLORS["CHAIN_OF_THOUGHT"]),
    ("🎭", "Persona /\nExpert Role", "Assigns specific expert identity that anchors vocabulary, tone, and perspective.",
     STRATEGY_COLORS["PERSONA_EXPERT_ROLE"]),
    ("📐", "Output Format\nSpecification", "Defines how output should be structured: length, layout, headers, tone, constraints.",
     STRATEGY_COLORS["OUTPUT_FORMAT_SPEC"]),
]

card_w = 2.9
gap = 0.15
start_x = 0.565
for i, (icon, name, desc, accent_color) in enumerate(strategies_1):
    x = start_x + i * (card_w + gap)
    # Card background
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, card_w, 4.0,
              fill=C["panelBg"], border=C["borderLine"])
    # Accent bar at top
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.1, card_w, 0.08, fill=accent_color)
    # Icon
    add_text_box(slide, x, 2.3, card_w, 0.4, icon, font_size=28, align="center")
    # Name
    add_text_box(slide, x + 0.15, 2.75, card_w - 0.3, 0.6, name,
                 font_size=14, bold=True, color=C["dark"], align="center")
    # Description
    add_text_box(slide, x + 0.15, 3.45, card_w - 0.3, 2.4, desc,
                 font_size=10.5, color=C["medGray"], align="center")


# ── Slide 7: Strategies Overview (5-8) ───────────────────────────────────────
slide = make_content_slide(
    "Eight prompting strategies — the AI's toolkit (2 of 2)",
    "THE STRATEGY ENGINE: STRATEGIES 5–8",
    "The remaining four strategies handle constraints, examples, decomposition, and tone — each filling a distinct need in the prompting craft.",
    "💡 Constraint Framing never appears alone — it's always additive. This teaches learners that constraints are modifiers, not drivers."
)

strategies_2 = [
    ("🚧", "Constraint\nFraming", "Scopes what the AI should NOT do — topics to avoid, length limits, risk prevention.",
     STRATEGY_COLORS["CONSTRAINT_FRAMING"]),
    ("📖", "Few-Shot\nExamples", "Includes concrete examples of desired output. AI pattern-matches structure, tone, and style.",
     STRATEGY_COLORS["FEW_SHOT_EXAMPLES"]),
    ("🔗", "Iterative\nDecomposition", "Breaks complex tasks into sequential sub-tasks, each building on the previous.",
     STRATEGY_COLORS["ITERATIVE_DECOMPOSITION"]),
    ("🎙️", "Tone & Voice\nSetting", "Specifies register and relational dynamic precisely: 'peer-to-peer, not evangelical.'",
     STRATEGY_COLORS["TONE_AND_VOICE"]),
]

for i, (icon, name, desc, accent_color) in enumerate(strategies_2):
    x = start_x + i * (card_w + gap)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, card_w, 4.0,
              fill=C["panelBg"], border=C["borderLine"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.1, card_w, 0.08, fill=accent_color)
    add_text_box(slide, x, 2.3, card_w, 0.4, icon, font_size=28, align="center")
    add_text_box(slide, x + 0.15, 2.75, card_w - 0.3, 0.6, name,
                 font_size=14, bold=True, color=C["dark"], align="center")
    add_text_box(slide, x + 0.15, 3.45, card_w - 0.3, 2.4, desc,
                 font_size=10.5, color=C["medGray"], align="center")


# ── Slide 8: Strategy Combination Rules ──────────────────────────────────────
slide = make_content_slide(
    "How strategies combine — the five rules",
    "STRATEGY COMBINATION LOGIC",
    "Not all strategies work together. These five rules govern which combinations produce effective prompts and which create noise.",
    "💡 Blueprint and Chain-of-Thought are rarely combined — choose based on whether the task needs structure or reasoning."
)

rules = [
    ("RULE 1", "Minimum 1, Maximum 4", "Every output uses at least one strategy. Four is the ceiling — beyond that, the prompt becomes overcrowded."),
    ("RULE 2", "Persona Almost Always Appears", "The only exception is purely mechanical tasks where identity is genuinely irrelevant."),
    ("RULE 3", "Blueprint ≠ Chain-of-Thought", "Choose the one that matches whether the task is about structure (Blueprint) or reasoning (CoT)."),
    ("RULE 4", "Constraints Are Additive", "Constraint Framing always appears alongside other strategies, never alone."),
    ("RULE 5", "Few-Shot ≠ Output Format", "Format Spec says how to structure. Few-Shot shows what good looks like. Combined = highly predictable."),
]

for i, (label, title, desc) in enumerate(rules):
    y = 2.1 + i * 0.82
    # Accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.565, y, 0.06, 0.72, fill=C["teal"])
    # Row background
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, y, 12.07, 0.72,
              fill=C["panelBg"], border=C["borderLine"])
    # Rule label
    add_text_box(slide, 0.85, y + 0.08, 1.2, 0.25, label,
                 font_size=9, bold=True, color=C["teal"])
    # Rule title
    add_text_box(slide, 0.85, y + 0.32, 3.5, 0.3, title,
                 font_size=12, bold=True, color=C["dark"])
    # Description
    add_text_box(slide, 4.5, y + 0.1, 8.1, 0.52, desc,
                 font_size=10.5, color=C["medGray"])


# ── Slide 9: Task Type → Strategy Map ────────────────────────────────────────
slide = make_content_slide(
    "Task type → strategy combination map",
    "STRATEGY COMBINATION LOGIC: TASK MAPPING",
    "Different task types call for different strategy combinations. This map guides the AI's selection logic and teaches learners through pattern recognition.",
    "💡 After 10 uses, learners notice: Persona appears almost always, Constraints never alone, Chain-of-Thought only for analytical tasks."
)

task_map = [
    ("Simple Communication", "Persona + Output Format", "2–3"),
    ("Analytical / Evaluative", "CoT + Constraints + Persona", "3–4"),
    ("Complex Deliverable", "Blueprint + Persona + Decomposition", "3–4"),
    ("Workshop Design", "Persona + CoT + Tone + Few-Shot", "3–4"),
    ("Template Creation", "Few-Shot + Output Format", "2–3"),
    ("Executive Comms", "Persona + Tone + Output Format", "3–4"),
    ("Creative Ideation", "Persona", "1–2"),
    ("Research Synthesis", "CoT + Decomposition + Constraints", "3–4"),
]

# Header row
headers = ["Task Type", "Primary Strategies", "Count"]
col_widths = [3.2, 7.0, 2.0]
col_x = [0.565, 3.765, 10.765]

for i, (header, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, 2.1, cw, 0.4,
              fill=C["dark"])
    add_text_box(slide, cx, 2.12, cw, 0.36, header,
                 font_size=10, bold=True, color=C["white"], align="center")

row_h = 0.47
for i, (task, strategies, count) in enumerate(task_map):
    y = 2.55 + i * row_h
    bg = C["white"] if i % 2 == 0 else C["panelBg"]
    for cx, cw in zip(col_x, col_widths):
        add_shape(slide, MSO_SHAPE.RECTANGLE, cx, y, cw, row_h, fill=bg, border=C["borderLine"])
    add_text_box(slide, col_x[0] + 0.1, y + 0.05, col_widths[0] - 0.2, row_h - 0.1,
                 task, font_size=10, bold=True, color=C["dark"])
    add_text_box(slide, col_x[1] + 0.1, y + 0.05, col_widths[1] - 0.2, row_h - 0.1,
                 strategies, font_size=10, color=C["medGray"])
    add_text_box(slide, col_x[2] + 0.1, y + 0.05, col_widths[2] - 0.2, row_h - 0.1,
                 count, font_size=10, bold=True, color=C["teal"], align="center")


# ── Slide 10: Section Divider — UI Design ────────────────────────────────────
make_divider_slide(
    "SECTION THREE",
    "UI Design Specification",
    "From input to output, every visual element is specified. The interface is designed to feel like a real work tool — not a learning exercise — while embedding educational moments in the output layer."
)


# ── Slide 11: Input Section ──────────────────────────────────────────────────
slide = make_content_slide(
    "Input section — single open textarea with example chips",
    "UI DESIGN: INPUT SECTION",
    "The input section removes all instructional scaffolding. A single open textarea with example chips keeps the interface focused and usable for real tasks.",
    "💡 Learning happens in the output, not the input. No questions, no guided forms — just describe your task."
)

# Wireframe-style layout
# Textarea wireframe
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 2.3, 9.3, 1.8,
          fill=C["white"], border=C["borderLine"])
add_text_box(slide, 2.2, 2.0, 5.0, 0.3,
             "Describe what you're trying to do", font_size=11, bold=True, color=C["dark"])

# Placeholder text in textarea
add_text_box(slide, 2.2, 2.45, 8.9, 0.5,
             "e.g., I need to write a project status update for my\nleadership team...",
             font_size=10.5, color=C["lightGray"])

# Character count
add_text_box(slide, 9.5, 3.75, 1.7, 0.2,
             "0 characters", font_size=9, color=C["lightGray"], align="right")

# Example chips
chip_labels = ["Write a stakeholder update", "Design an AI workshop", "Evaluate tools",
               "90-day onboarding plan", "Summarise research", "Draft a business case"]
chip_x = 2.0
for label in chip_labels:
    w = len(label) * 0.075 + 0.3
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, chip_x, 4.3, w, 0.35,
              fill=C["panelBg"], border=C["borderLine"])
    add_text_box(slide, chip_x + 0.1, 4.32, w - 0.2, 0.3,
                 label, font_size=9, color=C["bodyText"])
    chip_x += w + 0.12
    if chip_x > 10.5:
        break

# CTA Button
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 4.85, 2.2, 0.45,
          fill=C["teal"])
add_text_box(slide, 2.0, 4.87, 2.2, 0.4,
             "Build My Prompt →", font_size=12, bold=True, color=C["white"], align="center")

# Annotations
add_text_box(slide, 0.565, 5.6, 5.5, 0.3,
             "• Min-height 120px, auto-expands to 300px", font_size=9, color=C["medGray"])
add_text_box(slide, 0.565, 5.85, 5.5, 0.3,
             "• Focus state: border → #38B2AC (Teal)", font_size=9, color=C["medGray"])
add_text_box(slide, 6.5, 5.6, 5.5, 0.3,
             "• 5s client-side rate limit after submission", font_size=9, color=C["medGray"])
add_text_box(slide, 6.5, 5.85, 5.5, 0.3,
             "• Loading: 'Building your prompt...' with animated ellipsis", font_size=9, color=C["medGray"])


# ── Slide 12: Output Section — Two-Column Layout ─────────────────────────────
slide = make_content_slide(
    "Output section — two-column layout with strategy attribution",
    "UI DESIGN: OUTPUT SECTION",
    "The output appears with a slide-down + fade-in animation once the API responds. Left column (65%) shows the prompt; right column (35%) shows strategy cards.",
    "💡 The two-column layout makes the educational layer visible without interrupting the utility — the prompt is dominant, the learning is adjacent."
)

# Wireframe: two-column layout
# Left column (prompt)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.565, 2.2, 7.6, 3.9,
          fill=C["white"], border=C["borderLine"])
add_text_box(slide, 0.72, 2.25, 3.0, 0.2,
             "YOUR OPTIMISED PROMPT", font_size=9, bold=True, color=C["lightGray"])

# Original input preview
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 2.55, 7.3, 0.6,
          fill=C["panelBg"], border=C["borderLine"])
add_text_box(slide, 0.87, 2.58, 7.0, 0.55,
             "Your original input\n\"Write an update email about project progress...\"",
             font_size=9, color=C["medGray"])
add_text_box(slide, 0.72, 3.2, 7.3, 0.2,
             "Optimised to ↓", font_size=9, color=C["lightGray"], align="center")

# Prompt text area
add_text_box(slide, 0.87, 3.45, 7.0, 2.0,
             "You are a senior project manager writing a\nconcise progress update for a leadership\naudience who are time-poor and expect clarity\nover comprehensiveness.\n\nWrite a professional email that includes:\n• A 2-sentence summary of overall status\n• Top 2–3 milestones completed this week...",
             font_size=9.5, color=C["dark"])

# Action bar
action_labels = [("Copy Prompt", C["teal"], C["white"]),
                 ("← Try another", C["white"], C["dark"]),
                 ("Save as .txt", None, C["medGray"])]
ax = 0.72
for label, bg_color, text_color in action_labels:
    w = len(label) * 0.075 + 0.3
    if bg_color:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, ax, 5.7, w, 0.32,
                  fill=bg_color, border=C["borderLine"] if bg_color == C["white"] else None)
    add_text_box(slide, ax, 5.72, w, 0.28, label,
                 font_size=9, bold=(bg_color == C["teal"]), color=text_color, align="center")
    ax += w + 0.12

# Right column (strategy cards)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.4, 2.2, 4.37, 3.9,
          fill=C["panelBg"], border=C["borderLine"])
add_text_box(slide, 8.55, 2.25, 4.0, 0.2,
             "HOW THIS PROMPT WAS BUILT", font_size=9, bold=True, color=C["lightGray"])
add_text_box(slide, 8.55, 2.48, 4.0, 0.3,
             "2–4 strategies were combined. Click each to see why.",
             font_size=8, color=C["lightGray"])

# Mini strategy cards
mini_cards = [
    ("🎭", "Persona / Expert Role", STRATEGY_COLORS["PERSONA_EXPERT_ROLE"]),
    ("📐", "Output Format Specification", STRATEGY_COLORS["OUTPUT_FORMAT_SPEC"]),
]
for i, (icon, name, accent) in enumerate(mini_cards):
    y = 2.9 + i * 0.55
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.55, y, 4.07, 0.45,
              fill=C["white"], border=C["borderLine"])
    # Accent strip
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.55, y, 0.06, 0.45, fill=accent)
    add_text_box(slide, 8.72, y + 0.05, 3.7, 0.35,
                 f"{icon}  {name}", font_size=10, bold=True, color=C["dark"])
    add_text_box(slide, 12.15, y + 0.08, 0.4, 0.3,
                 "↓", font_size=10, color=C["lightGray"], align="center")

# Caveat banner
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.55, 4.15, 4.07, 1.0,
          fill=C["white"], border=C["borderLine"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 8.55, 4.15, 0.06, 1.0, fill=C["lightGray"])
add_text_box(slide, 8.72, 4.2, 3.8, 0.2,
             "There's no single perfect prompt.", font_size=9, bold=True, color=C["dark"])
add_text_box(slide, 8.72, 4.45, 3.8, 0.65,
             "This is an optimised starting point — not the only valid approach. The strategies above are the craft behind it.",
             font_size=8, color=C["medGray"])


# ── Slide 13: Strategy Cards Detail ──────────────────────────────────────────
slide = make_content_slide(
    "Strategy cards — collapsed and expanded states",
    "UI DESIGN: STRATEGY CARDS",
    "Each strategy used is rendered as a collapsible card with a left accent strip colour-coded to the strategy. Cards stack vertically with 10px gaps.",
    "💡 The 'Why this was used' rationale is generated dynamically per task — making it feel like a colleague's observation, not documentation."
)

# Collapsed card example
add_text_box(slide, 0.72, 2.15, 3.0, 0.25,
             "COLLAPSED STATE (DEFAULT)", font_size=9, bold=True, color=C["teal"])
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 2.45, 5.5, 0.55,
          fill=C["panelBg"], border=C["borderLine"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 2.45, 0.06, 0.55,
          fill=STRATEGY_COLORS["PERSONA_EXPERT_ROLE"])
add_text_box(slide, 0.92, 2.5, 4.5, 0.45,
             "🎭  Persona / Expert Role", font_size=12, bold=True, color=C["dark"])
add_text_box(slide, 5.72, 2.53, 0.4, 0.4,
             "↓", font_size=12, color=C["lightGray"], align="center")

# Expanded card example
add_text_box(slide, 0.72, 3.2, 3.0, 0.25,
             "EXPANDED STATE (ON CLICK)", font_size=9, bold=True, color=C["teal"])
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 3.5, 5.5, 2.7,
          fill=C["panelBg"], border=C["borderLine"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 3.5, 0.06, 2.7,
          fill=STRATEGY_COLORS["PERSONA_EXPERT_ROLE"])
add_text_box(slide, 0.92, 3.55, 4.5, 0.4,
             "🎭  Persona / Expert Role", font_size=12, bold=True, color=C["dark"])
add_text_box(slide, 5.72, 3.55, 0.4, 0.4,
             "↑", font_size=12, color=C["lightGray"], align="center")

# WHY section
add_text_box(slide, 0.92, 4.0, 5.0, 0.2,
             "WHY THIS WAS USED", font_size=8, bold=True, color=C["lightGray"])
add_text_box(slide, 0.92, 4.2, 5.0, 0.55,
             "A leadership audience expects authority and brevity — assigning a senior PM role shapes the voice before a word is written.",
             font_size=10, color=C["bodyText"])

# WHAT section
add_text_box(slide, 0.92, 4.85, 5.0, 0.2,
             "WHAT THIS STRATEGY DOES", font_size=8, bold=True, color=C["lightGray"])
add_text_box(slide, 0.92, 5.05, 5.0, 0.55,
             "Sets a specific expert identity that anchors the AI's register, vocabulary, and level of authority throughout the entire response.",
             font_size=10, color=C["medGray"])

# Styling specs on right side
specs = [
    "Card Styling:",
    "• Background: #F7FAFC",
    "• Border: 1px solid #E2E8F0",
    "• Border-radius: 10px",
    "• Padding: 12px 16px",
    "• Left accent: 4px, strategy colour",
    "",
    "Card Grid by Count:",
    "• 1–3 strategies: full-width stacked",
    "• 4 strategies: 2×2 grid",
    "• Mobile: always single column",
]
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.82, 2.45, 5.95, 3.75,
          fill=C["panelBg"], border=C["borderLine"])
add_multi_para_box(slide, 7.0, 2.55, 5.6, 3.5, specs,
                   font_size=10, color=C["medGray"], spacing=3)


# ── Slide 14: Section Divider — Pedagogy ─────────────────────────────────────
make_divider_slide(
    "SECTION FOUR",
    "The Pedagogy Layer",
    "Every design decision serves a learning outcome. This section defines the five principles that ensure the tool teaches prompting as a craft — not a formula — through repeated use rather than instruction."
)


# ── Slide 15: Learning Design Principles ─────────────────────────────────────
slide = make_content_slide(
    "Five learning design principles embedded in every interaction",
    "PEDAGOGY: LEARNING DESIGN PRINCIPLES",
    "The Prompt Playground's educational value is invisible by design — users feel like they're using a work tool while absorbing prompting principles through repeated exposure.",
    "💡 After three uses, a learner can say 'I need chain-of-thought here' rather than 'I need to add more detail.' That vocabulary is the learning."
)

principles = [
    ("P1", "Learning happens in the output, not the input",
     "No instructional scaffolding in the input stage. Intelligence lives in what comes back."),
    ("P2", "Strategies are named, not implied",
     "Every strategy is named explicitly — giving learners vocabulary they can carry forward."),
    ("P3", "Rationale is task-specific, not generic",
     "'Your task involves comparing options' feels like a colleague's observation, not documentation."),
    ("P4", "The caveat is the most important sentence",
     "'No single perfect prompt' positions prompting as a craft with judgment calls."),
    ("P5", "Repeated use compounds learning",
     "Patterns emerge through exposure: Persona always, Constraints never alone, CoT for analysis."),
]

for i, (label, title, desc) in enumerate(principles):
    y = 2.1 + i * 0.84
    # Number badge
    add_shape(slide, MSO_SHAPE.OVAL, 0.565, y + 0.05, 0.5, 0.5, fill=C["teal"])
    add_text_box(slide, 0.565, y + 0.1, 0.5, 0.4,
                 label, font_size=10, bold=True, color=C["white"], align="center")
    # Card
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, y, 11.57, 0.74,
              fill=C["panelBg"], border=C["borderLine"])
    add_text_box(slide, 1.35, y + 0.05, 11.3, 0.3,
                 title, font_size=12, bold=True, color=C["dark"])
    add_text_box(slide, 1.35, y + 0.38, 11.3, 0.3,
                 desc, font_size=10, color=C["medGray"])


# ── Slide 16: What Learners Absorb ───────────────────────────────────────────
slide = make_content_slide(
    "What learners absorb at each interaction moment",
    "PEDAGOGY: LEARNING JOURNEY",
    "Each moment in the user journey teaches something distinct — from articulating needs to recognising patterns across repeated use.",
    "💡 These patterns are not explained anywhere in the interface. They emerge through repeated use — the right model for building intuition."
)

moments = [
    ("Typing in the textarea", "Prompting starts with articulating what you actually need — most prompts fail here"),
    ("Receiving the prompt", "The gap between what I asked and what a good prompt looks like — and that it's closeable"),
    ("Opening strategy cards", "There are named techniques behind this, not magic — and I can learn them"),
    ("Reading 'Why this was used'", "Why this technique was right for this situation — begins building judgment"),
    ("Noticing pattern repetition", "Some techniques are universal (Persona), some situational, some risky to overuse"),
    ("Reading the caveat", "My job is not to memorise a formula but to develop a craft"),
]

for i, (moment, learning) in enumerate(moments):
    y = 2.1 + i * 0.7
    # Step number
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.565, y, 0.4, 0.55, fill=C["dark"])
    add_text_box(slide, 0.565, y + 0.1, 0.4, 0.35,
                 str(i + 1), font_size=12, bold=True, color=C["white"], align="center")
    # Moment
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.1, y, 3.8, 0.55,
              fill=C["teal"])
    add_text_box(slide, 1.2, y + 0.08, 3.6, 0.4,
                 moment, font_size=10, bold=True, color=C["white"])
    # Arrow
    add_text_box(slide, 5.0, y + 0.08, 0.3, 0.4,
                 "→", font_size=14, bold=True, color=C["teal"], align="center")
    # Learning
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.4, y, 7.37, 0.55,
              fill=C["panelBg"], border=C["borderLine"])
    add_text_box(slide, 5.55, y + 0.08, 7.1, 0.4,
                 learning, font_size=10, color=C["dark"])


# ── Slide 17: Section Divider — Technical ────────────────────────────────────
make_divider_slide(
    "SECTION FIVE",
    "Technical Specification",
    "The engineering details that power the experience — API configuration, state management, error handling, and responsive behaviour across breakpoints."
)


# ── Slide 18: API & Architecture ─────────────────────────────────────────────
slide = make_content_slide(
    "API specification and technical architecture",
    "TECHNICAL: API & ARCHITECTURE",
    "All API calls go through a backend proxy endpoint (Firebase Cloud Function via OpenRouter). The system prompt is stored server-side, never in the frontend.",
    "💡 The system prompt IS the intelligence layer — it contains the strategy selection logic, combination rules, and output format specification."
)

# Three-column layout
cols = [
    ("API Configuration", [
        "Model: claude-sonnet-4-20250514",
        "Max tokens: 1,500",
        "Endpoint: OpenRouter proxy",
        "Backend: Firebase Cloud Function",
        "Response: JSON (prompt + strategies)",
        "Timeout: 15 seconds",
    ]),
    ("State Management", [
        "userInput: string (textarea)",
        "isLoading: boolean (API call)",
        "result: object | null (response)",
        "expandedCards: Set<string>",
        "copySuccess: boolean",
        "originalInput: string (snapshot)",
    ]),
    ("Rate Limiting & Errors", [
        "Client-side: 5s after submission",
        "Parse failure → red error banner",
        "Timeout → amber warning notice",
        "API key missing → generic message",
        "Empty submit → inline validation",
        "Never expose API error details",
    ]),
]

col_w = 3.9
for i, (title, items) in enumerate(cols):
    x = 0.565 + i * (col_w + 0.15)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, col_w, 4.1,
              fill=C["panelBg"], border=C["borderLine"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.1, col_w, 0.08, fill=C["teal"])
    add_text_box(slide, x + 0.15, 2.25, col_w - 0.3, 0.3,
                 title, font_size=12, bold=True, color=C["dark"], align="center")
    add_multi_para_box(slide, x + 0.15, 2.65, col_w - 0.3, 3.3,
                       [f"• {item}" for item in items],
                       font_size=10, color=C["medGray"], spacing=6)


# ── Slide 19: Responsive Behaviour ───────────────────────────────────────────
slide = make_content_slide(
    "Responsive behaviour across three breakpoints",
    "TECHNICAL: RESPONSIVE DESIGN",
    "The layout adapts across desktop (1200px+), tablet (768–1199px), and mobile (<768px) with specific rules for the two-column output layout.",
    "💡 The two-column output is maintained on tablet to preserve the prompt-beside-strategies relationship — it only stacks on mobile."
)

breakpoints = [
    ("Desktop (1200px+)", C["teal"], [
        "Two-column output: 65% / 35%, gap 32px",
        "Max-width: 1100px, centred",
        "Textarea min-height: 120px",
        "Example chips: all visible",
        "4 strategy cards: 2×2 grid",
    ]),
    ("Tablet (768–1199px)", C["gold"], [
        "Two-column maintained: 60% / 40%, gap 24px",
        "Textarea min-height: 100px",
        "Chips: may require horizontal scroll",
        "Strategy cards: single-column stack",
        "Page padding: 24px",
    ]),
    ("Mobile (<768px)", C["orange"], [
        "Single column — stacked layout",
        "Prompt first, then strategy cards below",
        "CTA button: full width",
        "Action bar: stacks vertically, 8px gap",
        "Page padding: 16px",
    ]),
]

col_w = 3.9
for i, (label, accent, items) in enumerate(breakpoints):
    x = 0.565 + i * (col_w + 0.15)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, col_w, 4.1,
              fill=C["panelBg"], border=C["borderLine"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.1, col_w, 0.08, fill=accent)
    add_text_box(slide, x + 0.15, 2.25, col_w - 0.3, 0.35,
                 label, font_size=13, bold=True, color=C["dark"], align="center")
    add_multi_para_box(slide, x + 0.15, 2.7, col_w - 0.3, 3.3,
                       [f"• {item}" for item in items],
                       font_size=10, color=C["medGray"], spacing=8)


# ── Slide 20: Colour Palette & Typography ────────────────────────────────────
slide = make_content_slide(
    "Visual design specification — colours and typography",
    "TECHNICAL: DESIGN TOKENS",
    "Every colour and font weight in the Playground is specified. DM Sans is the typeface throughout with carefully defined size and weight scales.",
    "💡 Colour is not the sole differentiator for strategies — the strategy name and icon carry the same information, ensuring accessibility."
)

# Color swatches
colors_to_show = [
    ("Primary Text", "#1A202C", RGBColor(0x1A, 0x20, 0x2C)),
    ("Secondary Text", "#4A5568", RGBColor(0x4A, 0x55, 0x68)),
    ("Teal CTA", "#38B2AC", RGBColor(0x38, 0xB2, 0xAC)),
    ("Card BG", "#F7FAFC", RGBColor(0xF7, 0xFA, 0xFC)),
    ("Border", "#E2E8F0", RGBColor(0xE2, 0xE8, 0xF0)),
    ("Error", "#FC8181", RGBColor(0xFC, 0x81, 0x81)),
]

for i, (name, hex_val, rgb) in enumerate(colors_to_show):
    x = 0.565 + i * 2.05
    # Swatch
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.2, 1.85, 0.6, fill=rgb, border=C["borderLine"])
    add_text_box(slide, x, 2.85, 1.85, 0.2, name, font_size=8, bold=True, color=C["dark"], align="center")
    add_text_box(slide, x, 3.05, 1.85, 0.2, hex_val, font_size=8, color=C["medGray"], align="center")

# Strategy color swatches
strat_colors = [
    ("Blueprint", STRATEGY_COLORS["STRUCTURED_BLUEPRINT"]),
    ("Chain-of-Thought", STRATEGY_COLORS["CHAIN_OF_THOUGHT"]),
    ("Persona", STRATEGY_COLORS["PERSONA_EXPERT_ROLE"]),
    ("Output Format", STRATEGY_COLORS["OUTPUT_FORMAT_SPEC"]),
    ("Constraints", STRATEGY_COLORS["CONSTRAINT_FRAMING"]),
    ("Few-Shot", STRATEGY_COLORS["FEW_SHOT_EXAMPLES"]),
]
add_text_box(slide, 0.565, 3.45, 4.0, 0.25,
             "STRATEGY ACCENT COLOURS", font_size=9, bold=True, color=C["teal"])

for i, (name, rgb) in enumerate(strat_colors):
    x = 0.565 + i * 2.05
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 3.75, 1.85, 0.4, fill=rgb, border=C["borderLine"])
    add_text_box(slide, x, 4.2, 1.85, 0.2, name, font_size=8, bold=True, color=C["dark"], align="center")

# Typography table
add_text_box(slide, 0.565, 4.65, 4.0, 0.25,
             "TYPOGRAPHY — DM SANS", font_size=9, bold=True, color=C["teal"])

typo_data = [
    ("Page title", "700", "40–48px"),
    ("Input label", "600", "16px"),
    ("Textarea text", "400", "15px"),
    ("CTA button", "600", "15px"),
    ("Strategy card name", "600", "14px"),
    ("Card body", "400", "13px"),
]

for i, (element, weight, size) in enumerate(typo_data):
    y = 4.95 + i * 0.24
    add_text_box(slide, 0.72, y, 4.0, 0.22, element, font_size=9, color=C["dark"])
    add_text_box(slide, 5.0, y, 1.5, 0.22, weight, font_size=9, bold=True, color=C["teal"], align="center")
    add_text_box(slide, 6.8, y, 2.0, 0.22, size, font_size=9, color=C["medGray"])


# ── Slide 21: Worked Example ─────────────────────────────────────────────────
slide = make_content_slide(
    "Worked example — from user input to complete output",
    "WORKED EXAMPLE: SIMPLE COMMUNICATION TASK",
    "This shows the full flow for Example A from the PRD: a user asks for a project status update email and receives a Persona + Output Format optimised prompt.",
    "💡 Two strategies are sufficient for a simple email — over-engineering with 4 strategies would add noise, not signal."
)

# User input
add_text_box(slide, 0.72, 2.15, 4.0, 0.2,
             "USER INPUT", font_size=9, bold=True, color=C["teal"])
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 2.4, 11.8, 0.5,
          fill=C["panelBg"], border=C["borderLine"])
add_text_box(slide, 0.87, 2.45, 11.5, 0.4,
             "\"Write an update email about project progress to share with leadership\"",
             font_size=11, color=C["dark"])

# Strategies selected
add_text_box(slide, 0.72, 3.05, 4.0, 0.2,
             "STRATEGIES SELECTED", font_size=9, bold=True, color=C["teal"])

# Strategy pills
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 3.3, 2.5, 0.35,
          fill=STRATEGY_COLORS["PERSONA_EXPERT_ROLE"], border=C["borderLine"])
add_text_box(slide, 0.72, 3.32, 2.5, 0.3,
             "🎭  Persona / Expert Role", font_size=10, bold=True, color=C["dark"], align="center")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.4, 3.3, 3.0, 0.35,
          fill=STRATEGY_COLORS["OUTPUT_FORMAT_SPEC"], border=C["borderLine"])
add_text_box(slide, 3.4, 3.32, 3.0, 0.3,
             "📐  Output Format Specification", font_size=10, bold=True, color=C["white"], align="center")

# Generated prompt
add_text_box(slide, 0.72, 3.85, 4.0, 0.2,
             "GENERATED PROMPT", font_size=9, bold=True, color=C["teal"])
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 4.1, 7.5, 2.1,
          fill=C["white"], border=C["borderLine"])
add_text_box(slide, 0.87, 4.15, 7.2, 2.0,
             "You are a senior project manager writing a concise progress update for a leadership audience who are time-poor and expect clarity over comprehensiveness.\n\nWrite a professional email that includes:\n• A 2-sentence summary of overall project status\n• The top 2–3 milestones completed this week\n• Any active risks or blockers, each with a mitigation step\n• What decisions or support you need in the next two weeks\n\nKeep the email under 250 words. Use clear, scannable section headers.",
             font_size=9, color=C["dark"])

# Rationale cards on right
add_text_box(slide, 8.5, 3.85, 4.0, 0.2,
             "STRATEGY RATIONALE", font_size=9, bold=True, color=C["teal"])

# Persona rationale
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.5, 4.1, 4.27, 0.95,
          fill=C["panelBg"], border=C["borderLine"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 8.5, 4.1, 0.06, 0.95,
          fill=STRATEGY_COLORS["PERSONA_EXPERT_ROLE"])
add_text_box(slide, 8.65, 4.15, 4.0, 0.2,
             "🎭  Persona", font_size=10, bold=True, color=C["dark"])
add_text_box(slide, 8.65, 4.38, 4.0, 0.6,
             "A leadership audience expects authority and brevity — assigning a senior PM role shapes the voice before a word is written.",
             font_size=9, color=C["bodyText"])

# Output Format rationale
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.5, 5.15, 4.27, 0.95,
          fill=C["panelBg"], border=C["borderLine"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 8.5, 5.15, 0.06, 0.95,
          fill=STRATEGY_COLORS["OUTPUT_FORMAT_SPEC"])
add_text_box(slide, 8.65, 5.2, 4.0, 0.2,
             "📐  Output Format", font_size=10, bold=True, color=C["dark"])
add_text_box(slide, 8.65, 5.43, 4.0, 0.6,
             "Without format specification, a status email becomes a wall of prose. Defining the structure makes it scannable.",
             font_size=9, color=C["bodyText"])


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════

output_dir = "/Users/josephthomas/OXYGY-AI-Upskilling-Production-site-1/PRD"
output_path = os.path.join(output_dir, "Prompt_Playground_Redesign_PRD.pptx")
prs.save(output_path)
print(f"✅ Saved {prs.slides.__len__()} slides to: {output_path}")
